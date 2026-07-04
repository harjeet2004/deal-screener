#!/usr/bin/env python3
"""
deal_scout.py — MGA Ventures Weekly Deal-Sourcing Pipeline
===========================================================
Runs autonomously every Monday to surface one promising early-stage Indian
startup, produce a 9-slide PowerPoint deal memo, and email it.

All APIs used are on free tiers — zero cost per run:
  Serper.dev Search API   2 500 queries / month free (no credit card)
  Groq API (Llama 3.3)    14 400 requests / day free (no credit card)
  Gmail API               free with OAuth2

Run:  python deal_scout.py
"""

import argparse
import base64
import datetime
import io
import json
import os
import re
import sys
import time
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from google.auth.transport.requests import Request as GoogleRequest
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from groq import Groq
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Environment ───────────────────────────────────────────────────────────────
load_dotenv()

BASE_DIR        = Path(__file__).parent
CREDS_PATH      = BASE_DIR / "credentials.json"
TOKEN_PATH      = BASE_DIR / "token.json"

SERPER_API_KEY  = os.getenv("SERPER_API_KEY", "")
GROQ_API_KEY    = os.getenv("GROQ_API_KEY", "")
RECIPIENT_EMAIL = os.getenv("RECIPIENT_EMAIL", "")

GMAIL_SCOPES    = ["https://www.googleapis.com/auth/gmail.send"]

# ── Design palette ────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1A, 0x2B, 0x4A)
GOLD  = RGBColor(0xC9, 0xA8, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x2C, 0x2C, 0x2C)
GRAY  = RGBColor(0x88, 0x88, 0x88)
LTBG  = RGBColor(0xF4, 0xF6, 0xFA)
GREEN = RGBColor(0x27, 0x7A, 0x27)
AMBER = RGBColor(0xCC, 0x7A, 0x00)
RED_C = RGBColor(0xBB, 0x22, 0x22)

SECTORS = ["Consumer", "Fintech", "SaaS"]

# Sector-specific search query pairs (2 per sector).  Each string is inserted
# into the Serper query so results are tightly scoped to the chosen vertical.
_SECTOR_QUERIES: dict[str, list[str]] = {
    "consumer": [
        '"D2C" OR "consumer brand" OR "food tech" OR "quick commerce" OR "retail tech"',
        '"consumer internet" OR "FMCG startup" OR "beauty startup" OR "fashion startup"',
    ],
    "fintech": [
        '"payments" OR "digital lending" OR "insurtech" OR "neobank" OR "wealthtech"',
        '"BNPL" OR "embedded finance" OR "credit startup" OR "regtech" OR "paytech"',
    ],
    "saas": [
        '"B2B SaaS" OR "enterprise software" OR "SaaS startup" OR "vertical SaaS"',
        '"cloud platform" OR "API-first" OR "developer tools" OR "software startup"',
    ],
}

# Keywords that earn a bonus score when they appear in a result for a given sector.
_SECTOR_KEYWORDS: dict[str, list[str]] = {
    "consumer": ["d2c", "consumer brand", "fmcg", "food tech", "quick commerce",
                 "retail tech", "beauty", "fashion", "consumer internet"],
    "fintech":  ["fintech", "payments", "digital lending", "insurtech", "neobank",
                 "wealthtech", "bnpl", "embedded finance", "regtech"],
    "saas":     ["saas", "b2b saas", "enterprise software", "cloud platform",
                 "software", "api", "developer tools", "vertical saas"],
}

# Broader fallback queries used when primary queries yield no scored candidates.
# Intentionally simpler and wider — trading precision for recall.
_SECTOR_FALLBACK_QUERIES: dict[str, list[str]] = {
    "consumer": [
        'India consumer startup raises funding 2025 seed "Series A" crore',
        'India D2C brand startup funding 2025 investors backed',
    ],
    "fintech": [
        'India fintech startup raises funding 2025 seed "Series A" crore',
        'India payments startup OR lending startup funding 2025 investors',
    ],
    "saas": [
        'India SaaS startup raises funding 2025 seed "Series A" crore',
        'India software startup OR cloud startup funding 2025 investors backed',
    ],
}

# ─────────────────────────────────────────────────────────────────────────────
# 0. Sector rotation
# ─────────────────────────────────────────────────────────────────────────────

def sector_this_week() -> str:
    """Rotate Consumer → Fintech → SaaS by ISO week number."""
    week = datetime.date.today().isocalendar()[1]
    return SECTORS[week % 3]


# ─────────────────────────────────────────────────────────────────────────────
# 1. Serper.dev search helper
# ─────────────────────────────────────────────────────────────────────────────

def serper_search(query: str, num: int = 5) -> list:
    """Return up to `num` result dicts from Serper.dev (Google results, free tier)."""
    if not SERPER_API_KEY:
        print("[ERROR] SERPER_API_KEY not set in .env")
        return []
    try:
        r = requests.post(
            "https://google.serper.dev/search",
            headers={"X-API-KEY": SERPER_API_KEY, "Content-Type": "application/json"},
            json={"q": query, "num": min(num, 10), "gl": "in", "hl": "en"},
            timeout=15,
        )
        r.raise_for_status()
        items = r.json().get("organic", [])
        return [
            {
                "title":   i.get("title", ""),
                "snippet": i.get("snippet", ""),
                "link":    i.get("link", ""),
            }
            for i in items
        ]
    except Exception as exc:
        print(f"[WARN] Serper search failed for '{query[:50]}': {exc}")
        return []


# ─────────────────────────────────────────────────────────────────────────────
# 2. Page fetcher
# ─────────────────────────────────────────────────────────────────────────────

# Domains never fetched as page content (social, video, UGC, paywalled aggregators)
_FETCH_SKIP = {
    "facebook.com", "instagram.com", "twitter.com", "x.com", "t.co",
    "threads.net", "snapchat.com", "tiktok.com", "pinterest.com",
    "youtube.com", "youtu.be", "vimeo.com", "dailymotion.com",
    "reddit.com", "quora.com", "medium.com", "substack.com",
    "telegram.org", "t.me", "whatsapp.com",
    "linkedin.com", "glassdoor.com", "ambitionbox.com",
    "crunchbase.com", "tracxn.com", "growthlist.co",
    "wikipedia.org", "wikimedia.org",
}

# Domains blocked from discovery scoring and research corpus (blocklist)
_JUNK_DOMAINS = {
    "facebook.com", "instagram.com", "twitter.com", "x.com", "threads.net",
    "snapchat.com", "tiktok.com", "youtube.com", "reddit.com", "quora.com",
    "medium.com", "substack.com", "t.co", "linkedin.com",
    "growthlist.co", "tracxn.com", "crunchbase.com", "angellist.com",
    "f6s.com", "startupindia.gov.in",
    "wikipedia.org", "wikimedia.org",
    "thedigitalfifth.com", "statista.com", "ibef.org",
    "nasscom.in", "bcg.com", "mckinsey.com", "pwc.in", "kpmg.com",
    "deloitte.com", "ey.com", "gartner.com", "forrester.com",
    "inc42.com/resources", "yourstory.com/research",
}

# Allowlist: ONLY these domains are accepted into the research corpus.
# Any URL not matching this list is silently dropped before Groq sees it.
_CREDIBLE_DOMAINS = {
    # Indian startup / business news
    "inc42.com", "yourstory.com", "entrackr.com", "vccircle.com",
    "medianama.com", "themorningcontext.com", "the-ken.com",
    "businessinsider.in", "startupstorymedia.com",
    # Indian financial press
    "economictimes.indiatimes.com", "livemint.com", "mint.com",
    "business-standard.com", "financialexpress.com", "moneycontrol.com",
    "businesstoday.in", "cnbctv18.com", "zeebiz.com", "ndtv.com",
    "thehindu.com", "tribuneindia.com",
    # International financial / tech press
    "techcrunch.com", "reuters.com", "bloomberg.com", "wsj.com",
    "ft.com", "forbes.com", "fortune.com", "businesswire.com",
    "prnewswire.com", "globenewswire.com",
    # Deal data (read-only public pages)
    "dealstreetasia.com",
    # Professional network — snippets only (page fetch blocked by login wall)
    "linkedin.com",
}

def fetch_page(url: str, max_chars: int = 4000) -> str:
    """Fetch a web page and return cleaned plain text (first max_chars chars)."""
    if any(d in url for d in _FETCH_SKIP):
        return ""
    headers = {"User-Agent": "Mozilla/5.0 (compatible; MGADealScout/1.0)"}
    try:
        r = requests.get(url, headers=headers, timeout=12, allow_redirects=True)
        if r.status_code != 200:
            return ""
        soup = BeautifulSoup(r.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "form"]):
            tag.decompose()
        text = " ".join(soup.get_text(separator=" ").split())
        return text[:max_chars]
    except Exception:
        return ""


# ─────────────────────────────────────────────────────────────────────────────
# 3. Startup discovery
# ─────────────────────────────────────────────────────────────────────────────

def _score(title: str, snippet: str, url: str = "", sector: str = "") -> int:
    """Score how likely this result is an early-stage Indian startup (higher = better)."""
    text = (title + " " + snippet).lower()
    score = 0
    for kw in ["startup", "founded", "raises", "raised", "seed", "series a", "series b",
                "india", "crore", "lakh", "funding", "investors", "app", "platform",
                "inc42", "yourstory", "entrackr", "techcrunch"]:
        if kw in text:
            score += 2
    # Bonus for keywords that directly match the chosen sector
    for kw in _SECTOR_KEYWORDS.get(sector.lower(), []):
        if kw in text:
            score += 3
    # Extra bonus for Indian geography / currency signals
    for kw in ["bangalore", "bengaluru", "mumbai", "delhi", "hyderabad", "pune",
                "chennai", "noida", "gurugram", "inr", "rupee", "crore"]:
        if kw in text:
            score += 2
    for yr in ["2022", "2023", "2024", "2025"]:
        if yr in text:
            score += 1
    # Heavily penalise dead / distressed companies — not investable
    for kw in ["shuts down", "shut down", "shutdown", "closes down", "bankrupt",
               "winds down", "winding down", "lays off", "layoff", "job cut",
               "ceases operation", "acqui-hire", "employees seek clarity"]:
        if kw in text:
            score -= 15

    # Heavily penalise non-India HQ signals
    for kw in ["silicon valley", "san francisco", "new york", "new york-based", "nyc",
                "us-based", "us startup", "american startup", "bay area",
                "united states", "new zealand", "australia", "london-based",
                "singapore-based", "dubai-based", "toronto", "canada",
                "headquartered in us", "headquartered in new york",
                "founded in us", "founded in america"]:
        if kw in text:
            score -= 10
    # Penalise companies beyond Seed–Series B (outside MGA's stage mandate)
    for kw in ["series c", "series d", "series e", "series f", "series g",
               "pre-ipo", "ipo-bound", "ipo bound", "growth stage", "late stage",
               "late-stage", "pre ipo"]:
        if kw in text:
            score -= 12

    # Heavily penalise list/directory/aggregator pages — not actual startups
    for kw in ["unicorn", "ipo", "nyse", "bse", "nse", "nasdaq", "listed",
                "billion-dollar", "accelerator program", "vc fund", "venture capital firm",
                "wikipedia", "quora", "medium.com",
                "investment fund", "fund manager", "limited partners", "fund size",
                "ribbit", "sequoia", "accel", "tiger global", "softbank"]:
        if kw in text:
            score -= 4
    for kw in ["list of", "top 10", "top 50", "top 100", "3,928", "funded startups",
                "startup list", "directory", "database", "best startups", "startups to watch",
                "startups in india 2", "funded india startups", "ecosystem report",
                "fintech report", "saas report", "consumer report", "landscape report",
                "india fintech ecosystem", "india saas ecosystem", "india startup ecosystem",
                "report 2024", "report 2025", "report 2026",
                # Market-aggregate article patterns — NOT individual company news
                "funding reaches", "billion in 20", "investors adopt", "selective approach",
                "startup funding 2025", "startup ecosystem 2025", "market size",
                "billion raised", "total funding"]:
        if kw in text:
            score -= 8
    # Penalise aggregator/directory/report/tag URLs (not individual startup stories)
    for slug in ["/list", "/directory", "/database", "/top-", "/best-", "growthlist",
                 "tracxn.com/explore", "crunchbase.com/hub",
                 "/tag/", "/tags/", "/topic/", "/topics/", "/category/", "/author/",
                 "/explore/", "/search?", "/report/", "/monthly-funding",
                 "/weekly-funding", "/quarterly-funding", "/funding-report"]:
        if slug in url.lower():
            score -= 10
    return score


# Cities / phrases that indicate a non-India headquarters
_NON_INDIA_HQ_SIGNALS = {
    "new york", "new york-based", "san francisco", "silicon valley", "bay area",
    "los angeles", "seattle", "boston", "austin", "chicago", "toronto",
    "london", "amsterdam", "berlin", "paris", "singapore", "dubai",
    "us-based", "us startup", "american startup", "founded in us",
    "headquartered in us", "headquartered in new york", "new york headquarters",
}

def _is_india_based(title: str, snippet: str) -> bool:
    """Return False if the snippet/title clearly places HQ outside India."""
    text = (title + " " + snippet).lower()
    return not any(sig in text for sig in _NON_INDIA_HQ_SIGNALS)


_GENERIC_WORDS = {
    "india", "indian", "startup", "fintech", "saas", "consumer", "tech",
    "technology", "funding", "market", "ecosystem", "venture", "capital",
    "ventures", "partners", "fund", "investments",
}

def _extract_name(title: str) -> str:
    """
    Pull the startup name from a news headline.
    'TechCorp raises $2M seed — Inc42'  →  'TechCorp'
    Returns None if the title looks like a report/list, not a company name.
    """
    # Reject titles that are clearly reports, lists, aggregate stories, or dead companies
    reject_patterns = [
        "report", "ecosystem", "landscape", "funding trends",
        "top startup", "best startup", "list of", "funded india",
        "funding reaches", "billion in 20", "investors adopt",
        "market size", "startup funding 2025", "billion raised",
        "total funding", "selective approach", "roundup",
        "monthly funding", "weekly funding", "quarterly funding",
        "push august", "push september", "push october", "push november",
        "push december", "push january", "push february", "push march",
        "push april", "push may", "push june", "push july",
        "venture fund", "vc fund", "investment firm", "fund manager",
        # Dead / distressed company signals — not investable
        "shuts down", "shut down", "shutdown", "closes down", "ceases operation",
        "goes bankrupt", "files for bankruptcy", "files bankruptcy",
        "acquired by", "acqui-hire", "winds down", "winding down",
        "lays off", "mass layoff", "job cuts", "employees seek clarity",
    ]
    if any(p in title.lower() for p in reject_patterns):
        return None
    # Strip publication suffix
    for sep in [" — ", " - ", " | ", " : ", " · "]:
        if sep in title:
            title = title.split(sep)[0].strip()
    # Strip action verbs that follow the company name
    for suffix in [" raises", " raised", " raise ", " funding", " secures", " closes",
                   " gets ", " named", " launches", " acquires", " lands", " bags ",
                   " garners", " onboards", " partners"]:
        idx = title.lower().find(suffix)
        if idx > 3:
            title = title[:idx].strip()

    # Handle "Investor leads/backs round in [sector] startup CompanyName" pattern
    # e.g. "Bessemer leads seed round in SaaS startup Zenskar" → "Zenskar"
    m = re.search(r"\bstartup\s+([\w][\w\s\-\.]*?)\s*$", title, re.IGNORECASE)
    if m:
        candidate = m.group(1).strip()
        # Only accept if it looks like a proper name (not a generic word)
        if candidate and not all(w in _GENERIC_WORDS for w in candidate.lower().split()):
            return candidate

    # Strip common article prefixes like "Fintech startups Kiwi, CredRight"
    # that start with a sector/type label before the actual company names
    _PREFIX_RE = re.compile(
        r"^(?:fintech|saas|consumer|edtech|healthtech|agritech|deeptech|b2b|b2c)"
        r"\s+startup[s]?\s+", re.IGNORECASE
    )
    title = _PREFIX_RE.sub("", title).strip()

    # If multiple companies are listed (comma-separated), take the first one
    if "," in title:
        title = title.split(",")[0].strip()

    name = title.strip()
    if not name:
        return None

    words = name.split()

    # Company names are 1–4 words; longer strings are headlines or phrases
    if len(words) > 4:
        return None

    # Reject if any word is an article-action word — means this is a headline
    # phrase like "Fintech investors back scalable", not a company name
    _ARTICLE_ACTIONS = {
        "back", "backs", "backed", "backing",
        "fund", "funds", "funded",
        "invest", "invests", "invested",
        "support", "supports", "supported",
        "investors", "angel", "venture",
        "lead", "leads", "led", "leading",
        "scalable", "sustainable", "profitable",
    }
    if any(w.lower() in _ARTICLE_ACTIONS for w in words):
        return None

    # Reject if the extracted "name" is just generic sector/geography words
    # e.g. "India Startup", "Indian Fintech", "India SaaS"
    if all(w.lower() in _GENERIC_WORDS for w in words):
        return None

    return name


def discover_startup(sector: str, exclude: set | None = None) -> dict:
    """
    Run 2 sector-specific Serper queries scoped to Indian startup news sites,
    score all results, and return the best candidate that hasn't been tried yet.
    Raises SystemExit if no valid candidate is found.

    Args:
        sector:  One of "Consumer", "Fintech", "SaaS".
        exclude: Set of company names already tried this session — they are skipped.
    """
    exclude = {e.lower() for e in (exclude or set())}

    news_sites = (
        "site:inc42.com OR site:yourstory.com OR site:entrackr.com "
        "OR site:economictimes.indiatimes.com OR site:vccircle.com "
        "OR site:medianama.com OR site:themorningcontext.com"
    )

    sector_key = sector.lower()
    q_terms = _SECTOR_QUERIES.get(sector_key, [f'"{sector.lower()} startup"', ""])
    primary_queries = [
        f'{news_sites} India {q_terms[0]} startup raises crore 2025 -report -roundup -list',
        f'{news_sites} India {sector.lower()} {q_terms[1]} "seed" OR "Series A" OR "Series B" 2025 -report -list',
    ]
    fallback_queries = _SECTOR_FALLBACK_QUERIES.get(sector_key, [
        f'India {sector.lower()} startup raises seed funding 2025',
    ])

    def _run_queries(qs: list[str]) -> list:
        results = []
        for q in qs:
            results.extend(serper_search(q, num=8))
            time.sleep(0.6)
        return results

    def _score_results(raw: list) -> list:
        seen, out = set(), []
        for r in raw:
            url    = r["link"]
            domain = re.sub(r"https?://(?:www\.)?", "", url).split("/")[0]
            if url in seen or domain in _JUNK_DOMAINS:
                continue
            seen.add(url)
            s = _score(r["title"], r["snippet"], url, sector=sector)
            if s > 0:
                out.append((s, r))
        return out

    # Try primary queries first; fall back to broader queries if nothing scored
    all_results = _run_queries(primary_queries)
    scored = _score_results(all_results)

    if not scored:
        print("[Discovery] Primary queries yielded no candidates — trying fallback queries…")
        fallback_results = _run_queries(fallback_queries)
        scored = _score_results(fallback_results)

    if not scored:
        print("[ERROR] No suitable startup candidate found even with fallback queries.")
        sys.exit(1)

    scored.sort(key=lambda x: x[0], reverse=True)
    for _, best in scored:
        name = _extract_name(best["title"])
        if not name:
            continue
        if name.lower() in exclude:
            print(f"[Discovery] Skipping '{name}' — already tried this session")
            continue
        if not _is_india_based(best["title"], best["snippet"]):
            print(f"[Discovery] Skipping '{name}' — HQ appears to be outside India")
            continue
        print(f"[Discovery] Candidate: {name}  ({best['link']})")
        return {"name": name, "url": best["link"], "snippet": best["snippet"]}

    print("[ERROR] No new India-based startup candidate found (all results already tried).")
    sys.exit(1)


# ─────────────────────────────────────────────────────────────────────────────
# 4. Deep research
# ─────────────────────────────────────────────────────────────────────────────

def deep_research(company: str, sector: str, discovery_url: str) -> dict:
    """
    Run 4 targeted CSE queries + fetch pages for the candidate company.
    Returns {corpus: str, sources: list[str]}.
    Uses 4 CSE queries (plus the 2 from discovery = 6 total per run,
    well within the 100/day free limit).
    """
    queries = [
        f'"{company}" India startup product service what it does',
        f'"{company}" India founders CEO CTO team background',
        f'"{company}" India funding raised investors crore rupee round',
        f'"{company}" India revenue users traction growth metrics',
    ]

    corpus, sources = [], []

    # Always start with the discovery page
    print(f"[Research] Fetching discovery page …")
    page = fetch_page(discovery_url)
    if page:
        corpus.append(f"SOURCE: {discovery_url}\n{page}")
    if discovery_url not in sources:
        sources.append(discovery_url)

    for q in queries:
        print(f"[Research] Searching: {q[:65]} …")
        results = serper_search(q, num=3)
        time.sleep(0.6)

        for idx, r in enumerate(results):
            link = r["link"]
            domain = re.sub(r"https?://(?:www\.)?", "", link).split("/")[0]
            # Allowlist: only credible financial / startup news outlets enter the corpus
            if not any(domain == c or domain.endswith("." + c) for c in _CREDIBLE_DOMAINS):
                continue
            entry = (
                f"SOURCE: {link}\n"
                f"TITLE: {r['title']}\n"
                f"SNIPPET: {r['snippet']}"
            )
            corpus.append(entry)
            if link not in sources:
                sources.append(link)

            # Fetch the top result for each query only
            if idx == 0:
                txt = fetch_page(link)
                if txt:
                    corpus.append(f"PAGE CONTENT ({link}):\n{txt}")

    print(f"[Research] Collected {len(sources)} sources.")
    return {"corpus": "\n\n---\n\n".join(corpus), "sources": sources}


# ─────────────────────────────────────────────────────────────────────────────
# 5. Gemini synthesis
# ─────────────────────────────────────────────────────────────────────────────

_PROMPT = """\
You are a deal analyst at MGA Ventures, a Mumbai family office investing in \
early-stage Indian startups (Seed to Series B, Consumer / Fintech / SaaS, \
ticket up to $2M USD).

Below is raw web research on a company called "{company}" in the {sector} \
sector. Extract ONLY facts explicitly stated in the text.

CRITICAL RULES — no exceptions:
• NEVER invent, estimate, or extrapolate any number, name, date, or fact.
• If information is absent from the text write exactly: Not publicly disclosed.
• Every fact must be traceable to a URL in the source text.
• Only cite credible, verifiable sources: news outlets (Inc42, YourStory, Entrackr, \
Economic Times, Mint, Business Standard, VCCircle, TechCrunch, Reuters, Bloomberg), \
LinkedIn company/founder pages, company websites, or official press releases. \
NEVER cite Facebook, Instagram, Twitter/X, Reddit, YouTube, Quora, or anonymous/UGC pages.
• Be concise and factual. No filler phrases.

Return ONLY valid JSON — no markdown fences, no commentary — matching this \
exact structure:
{{
  "company_name": "...",
  "tagline": "one sentence tagline, 12 words or fewer",
  "hq": "City, India   or   Not publicly disclosed.",
  "founded_year": "YYYY   or   Not publicly disclosed.",
  "website": "URL   or   Not publicly disclosed.",
  "stage": "Seed / Pre-Series A / Series A / Series B / Not publicly disclosed.",
  "sector": "{sector}",
  "problem": "2-3 sentences from sources only.",
  "solution": "2-3 sentences from sources only.",
  "business_model": "How it makes money, from sources only.",
  "products": [
    {{"name": "Product or Feature Name", "description": "What it does — 1-2 sentences from sources only."}}
  ],
  "founders": [
    {{"name": "...", "title": "...", "background": "one sentence from sources", "linkedin_url": "LinkedIn profile URL if found in sources, else ''"}}
  ],
  "market_size": {{
    "tam": {{"value": "Total Addressable Market — ₹X Cr or $XB from sources, or 'Not publicly disclosed.'", "source_url": "URL that contained this figure, or ''"}},
    "sam": {{"value": "Serviceable Addressable Market — ₹X Cr or $XB from sources, or 'Not publicly disclosed.'", "source_url": ""}},
    "som": {{"value": "Serviceable Obtainable Market — ₹X Cr or $XB from sources, or 'Not publicly disclosed.'", "source_url": ""}},
    "context": "2-3 sentences of broader market context from sources."
  }},
  "funding_rounds": [
    {{"round": "Seed/Pre-A/Series A/etc.", "amount": "₹X Cr or $XM", "investors": "names", "date": "Mon YYYY", "source_url": "URL from which this round was confirmed — required, do not leave blank if the round is mentioned in any source"}}
  ],
  "traction_metrics": [
    {{"metric": "MAU / ARR / GMV / etc.", "value": "exact figure from sources", "source_url": "URL"}}
  ],
  "key_differentiators": [
    "Key competitive differentiator or moat found in sources — must attribute to an outlet e.g. 'proprietary AI underwriting model (Inc42)'",
    "Second differentiator from sources"
  ],
  "thesis_fit_stage": "Fit / Partial fit / Out of scope — one sentence explanation.",
  "thesis_fit_sector": "Fit / Partial fit / Out of scope — one sentence.",
  "thesis_fit_ticket": "Fit / Partial fit / Out of scope — one sentence.",
  "red_flags": ["specific red flag — skip this field if you have none"],
  "investment_highlights": [
    "Most compelling investment reason — quote specific data e.g. '3× YoY revenue growth (Inc42)'",
    "Second highlight from sources",
    "Third highlight"
  ],
  "competitive_landscape": [
    {{"competitor": "Direct competitor name or 'Manual/Traditional processes'", "differentiator": "Why this company is better or different — from sources only"}}
  ],
  "use_of_funds": "How the company plans to deploy raised capital from sources, or 'Not publicly disclosed.'",
  "partnerships": ["Key enterprise customer, distribution partner or platform from sources — omit entire list if none found"],
  "awards_recognition": ["Award, accelerator batch, or industry recognition with year — omit if none found"],
  "team_strength": "2 sentences on founding team quality, domain expertise, and execution track record based solely on sources.",
  "recent_news": [
    {{"headline": "Latest news headline verbatim or closely paraphrased from source", "date": "Mon YYYY or ''", "source_url": "URL"}}
  ],
  "email_summary": "Exactly 3 sentences: what the company does, its traction, and why it fits MGA's thesis.",
  "sources": [
    CRITICAL: This array MUST contain one entry for EVERY number, amount, name, date, and
    statistic you placed anywhere in this JSON. That means:
      • One entry per funding_round (claim = the round amount + investor names)
      • One entry per traction_metric (claim = the exact figure you listed)
      • One entry per TAM / SAM / SOM figure
      • One entry per founder background sentence
      • One entry per key_differentiator or investment_highlight that names a specific fact
    Do NOT leave this sparse. If you listed a number above, there must be a source here for it.
    {{
      "url": "full URL of the article where this specific fact appears",
      "outlet": "Publication name e.g. Inc42, Economic Times, YourStory",
      "date": "Mon YYYY or ''",
      "claim": "The SPECIFIC fact — must quote the actual number, name, or date.  E.g.: 'raised Rs 100 Cr Series A from BEENEXT and Fireside Ventures (May 2025)' or '5 lakh monthly active users as of Q4 2024'"
    }}
  ]
}}

RAW RESEARCH TEXT:
{corpus}
"""


def synthesise(company: str, sector: str, research: dict) -> dict:
    """Call Groq (Llama 3) to turn raw research into a structured deal dict."""
    client = Groq(api_key=GROQ_API_KEY)

    prompt = _PROMPT.format(
        company=company,
        sector=sector,
        corpus=research["corpus"][:24000],
    )

    print("[Groq] Synthesising …")
    for attempt in range(3):
        try:
            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.1,
            )
            raw = resp.choices[0].message.content.strip()
            raw = re.sub(r"^```(?:json)?\s*", "", raw)
            raw = re.sub(r"\s*```$",           "", raw)
            data = json.loads(raw)
            print("[Groq] Done.")
            return data
        except json.JSONDecodeError as exc:
            print(f"[WARN] Groq returned invalid JSON (attempt {attempt + 1}): {exc}")
            if attempt < 2:
                time.sleep(5)
        except Exception as exc:
            print(f"[WARN] Groq error (attempt {attempt + 1}): {exc}")
            if attempt < 2:
                time.sleep(10)

    # Graceful fallback — deck still builds, just with minimal content
    print("[ERROR] Gemini synthesis failed after 3 attempts — using fallback skeleton.")
    return {
        "company_name": company, "tagline": "Not publicly disclosed.",
        "hq": "India", "founded_year": "Not publicly disclosed.",
        "website": "Not publicly disclosed.", "stage": "Not publicly disclosed.",
        "sector": sector, "problem": "Not publicly disclosed.",
        "solution": "Not publicly disclosed.", "business_model": "Not publicly disclosed.",
        "products": [], "founders": [],
        "market_size": {"tam": {"value": "Not publicly disclosed.", "source_url": ""},
                        "sam": {"value": "Not publicly disclosed.", "source_url": ""},
                        "som": {"value": "Not publicly disclosed.", "source_url": ""},
                        "context": "Not publicly disclosed."},
        "funding_rounds": [], "traction_metrics": [], "key_differentiators": [],
        "thesis_fit_stage": "Not publicly disclosed.",
        "thesis_fit_sector": "Not publicly disclosed.",
        "thesis_fit_ticket": "Not publicly disclosed.",
        "red_flags": [],
        "investment_highlights": [], "competitive_landscape": [],
        "use_of_funds": "Not publicly disclosed.", "partnerships": [],
        "awards_recognition": [], "team_strength": "Not publicly disclosed.",
        "recent_news": [],
        "email_summary": f"{company} — Gemini synthesis failed; review raw sources manually.",
        "sources": [{"url": u, "supports": "research source"} for u in research["sources"]],
    }


# ─────────────────────────────────────────────────────────────────────────────
# 6. PowerPoint builder
# ─────────────────────────────────────────────────────────────────────────────

def _blank_layout(prs):
    """Return the Blank slide layout (robust across different pptx themes)."""
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", "blank slide"):
            return layout
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def _add_slide(prs):
    return prs.slides.add_slide(_blank_layout(prs))


def _txb(slide, text, left, top, width, height,
         size=13, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    """Add a textbox. Coordinates and dimensions in inches."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = str(text)
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = color
    return txb


_LINK_BLUE = RGBColor(0x1A, 0x6B, 0xC4)


def _txb_link(slide, text, href, left, top, width, height,
              size=9, bold=False):
    """Add a textbox whose text renders as a clickable hyperlink (blue, underlined)."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.underline = True
    run.font.color.rgb = _LINK_BLUE
    if href:
        run.hyperlink.address = href
    return txb


def _header_bar(slide, title: str, subtitle: str = ""):
    """Navy header bar across the top of a content slide."""
    bar_h = 1.1 if subtitle else 0.95
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    _txb(slide, title, 0.35, 0.08, 11.5, 0.7,
         size=26, bold=True, color=WHITE)
    if subtitle:
        _txb(slide, subtitle, 0.35, 0.72, 12.5, 0.34,
             size=14, color=RGBColor(0xB8, 0xC4, 0xD8))


def _kv_table(slide, rows, left=0.4, top=1.15, col_widths=(2.4, 5.8)):
    """Render (key, value) pairs as a clean two-column table."""
    n   = len(rows)
    w   = sum(col_widths)
    tbl = slide.shapes.add_table(
        n, 2, Inches(left), Inches(top), Inches(w), Inches(n * 0.56)
    ).table
    tbl.columns[0].width = Inches(col_widths[0])
    tbl.columns[1].width = Inches(col_widths[1])
    for i, (k, v) in enumerate(rows):
        c0 = tbl.cell(i, 0)
        c0.text = str(k)
        c0.text_frame.paragraphs[0].font.bold      = True
        c0.text_frame.paragraphs[0].font.size      = Pt(15)
        c0.text_frame.paragraphs[0].font.color.rgb = NAVY
        c0.fill.solid()
        c0.fill.fore_color.rgb = LTBG

        c1 = tbl.cell(i, 1)
        c1.text = str(v)
        c1.text_frame.paragraphs[0].font.size      = Pt(15)
        c1.text_frame.paragraphs[0].font.color.rgb = DARK
        if i % 2 == 1:
            c1.fill.solid()
            c1.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFD)


def _bullet_list(slide, items, left, top, width, height, size=13):
    """Add a bulleted list."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size      = Pt(size)
        p.font.color.rgb = DARK
        p.space_before   = Pt(5)


def _traction_chart(metrics: list) -> bytes:
    """
    Generate a horizontal bar chart from traction metrics that have numeric values.
    Returns PNG bytes, or empty bytes if no numeric data is found.
    """
    numeric = []
    for m in metrics:
        val_str = str(m.get("value", ""))
        # Extract the first number in the value string
        match = re.search(r"[\d,]+(?:\.\d+)?", val_str.replace(",", ""))
        if match:
            try:
                numeric.append((m.get("metric", ""), float(match.group()), val_str))
            except ValueError:
                pass

    if not numeric:
        return b""

    labels = [f"{m}\n({v})" for m, _, v in numeric[:6]]
    values = [n for _, n, _ in numeric[:6]]

    fig, ax = plt.subplots(figsize=(6.5, max(2.2, len(labels) * 0.7)))
    ax.barh(labels, values, color="#1A2B4A", edgecolor="white", height=0.55)
    ax.set_xlabel("Value", fontsize=10)
    ax.tick_params(axis="y", labelsize=9)
    ax.spines[["top", "right"]].set_visible(False)
    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _funding_chart(rounds: list) -> bytes:
    """Bar chart of funding rounds. Falls back to timeline if amounts aren't parseable."""
    if not rounds:
        return b""

    labels, values = [], []
    for r in rounds:
        amt = r.get("amount", "")
        # Require at least one digit before an optional decimal — avoids matching bare "."
        m = re.search(r"\d+\.?\d*", amt.replace(",", ""))
        if not m:
            continue
        try:
            val = float(m.group())
        except ValueError:
            continue
            # Rough normalisation to Cr: $1M ~ 85 Cr
            if "$" in amt and any(x in amt.upper() for x in ["M", "MN", "MILLION"]):
                val *= 85
            elif "$" in amt and "K" in amt.upper():
                val *= 0.085
            labels.append(r.get("round", "Round"))
            values.append(val)

    if not values:
        # Show a text-based timeline instead
        return _round_timeline(rounds)

    fig, ax = plt.subplots(figsize=(5.5, max(2.8, len(labels) * 1.0)))
    bars = ax.bar(range(len(labels)), values, color="#C9A84C",
                  edgecolor="#1A2B4A", linewidth=1.5, width=0.5)
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylabel("Amount (Cr)", fontsize=9, color="#1A2B4A")
    ax.set_title("Funding History", fontsize=11, fontweight="bold", color="#1A2B4A", pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(axis="both", labelsize=9)
    top = max(values)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + top * 0.02,
                f"{val:.0f}Cr", ha="center", va="bottom", fontsize=8,
                color="#1A2B4A", fontweight="bold")
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _round_timeline(rounds: list) -> bytes:
    """Horizontal timeline for funding rounds when no numeric amounts are available."""
    items = rounds[:5]
    labels  = [r.get("round", "?") for r in items]
    amounts = [r.get("amount", "?") for r in items]
    dates   = [r.get("date", "") for r in items]
    n = len(items)

    fig, ax = plt.subplots(figsize=(5.5, 2.2))
    ax.set_xlim(-0.3, n - 0.7)
    ax.set_ylim(-1.2, 1.5)
    ax.axis("off")
    ax.axhline(y=0, color="#1A2B4A", linewidth=2, xmin=0.03, xmax=0.97)
    for i, (label, amount, date) in enumerate(zip(labels, amounts, dates)):
        ax.plot(i, 0, "o", color="#C9A84C", markersize=18, zorder=3)
        ax.text(i, 0, label, ha="center", va="center",
                fontsize=7, fontweight="bold", color="#1A2B4A", zorder=4)
        ax.text(i, 0.55, amount, ha="center", va="bottom", fontsize=8,
                color="#2C2C2C", fontweight="bold")
        ax.text(i, -0.45, date, ha="center", va="top", fontsize=7, color="#888888")
    ax.set_title("Funding Timeline", fontsize=10, fontweight="bold", color="#1A2B4A", pad=4)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _market_donut(market_size: dict) -> bytes:
    """
    Concentric-circle market chart from structured TAM/SAM/SOM data.
    Only draws circles for tiers that have real numeric values from sources.
    """
    tiers = [
        ("TAM", market_size.get("tam", {}).get("value", "")),
        ("SAM", market_size.get("sam", {}).get("value", "")),
        ("SOM", market_size.get("som", {}).get("value", "")),
    ]
    colors = ["#1A2B4A", "#2E5FA3", "#C9A84C"]

    parsed = []
    for lbl, val_str in tiers:
        if not val_str or "not publicly" in val_str.lower():
            continue
        # Extract numeric value in $B equivalent
        b = re.search(r"\$\s*([\d.]+)\s*(?:billion|bn|b)\b", val_str, re.IGNORECASE)
        m = re.search(r"\$\s*([\d.]+)\s*(?:million|mn|m)\b",  val_str, re.IGNORECASE)
        c = re.search(r"(?:rs\.?|inr|₹)?\s*([\d,.]+)\s*(?:cr|crore)",  val_str, re.IGNORECASE)
        if b:
            parsed.append((lbl, float(b.group(1)), val_str.split("—")[0].strip()))
        elif m:
            parsed.append((lbl, float(m.group(1)) / 1000, val_str.split("—")[0].strip()))
        elif c:
            parsed.append((lbl, float(c.group(1).replace(",","")) / 7000, val_str.split("—")[0].strip()))

    if not parsed:
        return b""   # don't draw a chart with invented numbers

    fig, ax = plt.subplots(figsize=(4.5, 4.0))
    for i, (lbl, _, display) in enumerate(parsed):
        frac = 0.88 - i * 0.24
        col  = colors[i % len(colors)]
        circle = plt.Circle((0.5, 0.5), frac / 2, color=col,
                             alpha=0.85 - i * 0.1, transform=ax.transAxes, zorder=3 - i)
        ax.add_patch(circle)
        offset = 0.28 - i * 0.20
        ax.text(0.5, 0.5 + offset, f"{lbl}: {display}",
                ha="center", va="center", fontsize=9, fontweight="bold",
                color="white", transform=ax.transAxes, zorder=5)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_title("Market Size (Sourced)", fontsize=11, fontweight="bold",
                 color="#1A2B4A", pad=6)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _fit_score_val(text: str) -> int:
    t = text.lower()
    if "out of scope" in t or "not publicly" in t:
        return 1
    if "partial" in t:
        return 2
    return 3


def _fit_color_hex(text: str) -> str:
    v = _fit_score_val(text)
    return {1: "#BB2222", 2: "#CC7A00", 3: "#277A27"}[v]


def _thesis_gauge(criteria: list) -> bytes:
    """Horizontal bar scorecard showing thesis fit per criterion."""
    labels = [c[0] for c in criteria]
    values = [_fit_score_val(c[1]) for c in criteria]
    colors = [_fit_color_hex(c[1]) for c in criteria]

    fig, ax = plt.subplots(figsize=(4.2, 2.4))
    bars = ax.barh(labels, values, color=colors, height=0.45, edgecolor="white", linewidth=1)
    ax.set_xlim(0, 3.4)
    ax.set_xticks([1, 2, 3])
    ax.set_xticklabels(["Out of\nScope", "Partial\nFit", "Fit"], fontsize=8, color="#444")
    ax.axvline(x=1, color="#ddd", linewidth=0.8, linestyle="--")
    ax.axvline(x=2, color="#ddd", linewidth=0.8, linestyle="--")
    ax.spines[["top", "right", "bottom"]].set_visible(False)
    ax.tick_params(axis="y", labelsize=9)
    ax.set_title("Thesis Fit Scorecard", fontsize=10, fontweight="bold",
                 color="#1A2B4A", pad=6)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _fit_color(text: str) -> RGBColor:
    t = text.lower()
    if "out of scope" in t or "not publicly" in t:
        return RED_C
    if "partial" in t:
        return AMBER
    return GREEN


def _fetch_logo(website: str) -> bytes:
    """
    Fetch company logo PNG via Clearbit Logo API (free, no auth needed).
    Returns raw PNG bytes on success, b'' if unavailable.
    """
    if not website or "not publicly" in website.lower():
        return b""
    try:
        domain = re.sub(r"https?://(?:www\.)?", "", website).split("/")[0].rstrip("/")
        if not domain or "." not in domain:
            return b""
        r = requests.get(
            f"https://logo.clearbit.com/{domain}",
            timeout=8,
            headers={"User-Agent": "Mozilla/5.0"},
        )
        if r.status_code == 200 and len(r.content) > 300:
            print(f"[Logo] Fetched logo for {domain}")
            return r.content
    except Exception as exc:
        print(f"[Logo] Could not fetch logo: {exc}")
    return b""


def _place_logo(slide, logo_bytes: bytes, left: float, top: float,
                size: float = 1.4, dark_bg: bool = False):
    """Embed a logo on a slide with an optional white backing card."""
    if not logo_bytes:
        return
    try:
        if dark_bg:
            # White pill so logo is visible on dark backgrounds
            bg = slide.shapes.add_shape(
                1, Inches(left - 0.12), Inches(top - 0.1),
                Inches(size + 0.24), Inches(size * 0.58 + 0.2),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = WHITE
            bg.line.fill.background()
        slide.shapes.add_picture(
            io.BytesIO(logo_bytes),
            Inches(left), Inches(top),
            width=Inches(size),
        )
    except Exception:
        pass


def build_deck(d: dict, today_str: str, logo_bytes: bytes = b"") -> Path:
    """Build the 9-slide PowerPoint and return its Path."""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    company = d.get("company_name", "Unknown Startup")

    # ── Slide 1 · Title ───────────────────────────────────────────────────────
    s1 = _add_slide(prs)
    s1.background.fill.solid()
    s1.background.fill.fore_color.rgb = NAVY

    _txb(s1, "MGA VENTURES",
         0.5, 0.5, 12.33, 0.55, size=14, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
    _txb(s1, "DEAL SCREENING MEMO",
         0.5, 1.0, 12.33, 0.45, size=11,
         color=RGBColor(0xA8, 0xB8, 0xCC), align=PP_ALIGN.CENTER)
    _txb(s1, company.upper(),
         0.5, 1.7, 12.33, 1.8, size=52, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _txb(s1, d.get("tagline", ""),
         0.5, 3.55, 12.33, 0.65, size=18, italic=True,
         color=GOLD, align=PP_ALIGN.CENTER)
    sector_stage = "  ·  ".join(filter(None, [d.get("sector", ""), d.get("stage", "")]))
    _txb(s1, sector_stage,
         0.5, 4.25, 12.33, 0.45, size=13,
         color=RGBColor(0xB0, 0xBC, 0xCC), align=PP_ALIGN.CENTER)

    # One-line company summary (first sentence of email_summary)
    summary_full = d.get("email_summary", "")
    if summary_full:
        first_sent = (summary_full.split(".")[0].strip() + ".") if "." in summary_full else summary_full
        _txb(s1, first_sent, 1.0, 4.88, 11.33, 0.82,
             size=15, italic=True, color=RGBColor(0xC0, 0xCC, 0xDD),
             align=PP_ALIGN.CENTER)

    # Four info chips: Founded | HQ | Stage | Funding
    total_funding = ""
    for fr in (d.get("funding_rounds") or []):
        if fr.get("amount"):
            total_funding = fr["amount"]
    chips1 = [
        ("FOUNDED",  d.get("founded_year", "N/A")),
        ("HQ",       d.get("hq", "India")),
        ("STAGE",    d.get("stage", "N/A")),
        ("FUNDING",  total_funding or d.get("sector", "")),
    ]
    chip1_w = 2.9
    for j, (clbl, cval) in enumerate(chips1):
        cxj = 0.47 + j * (chip1_w + 0.25)
        cb = s1.shapes.add_shape(1, Inches(cxj), Inches(5.88),
                                 Inches(chip1_w), Inches(0.82))
        cb.fill.solid()
        cb.fill.fore_color.rgb = RGBColor(0x22, 0x38, 0x5F)
        cb.line.color.rgb = GOLD
        cb.line.width = Pt(0.75)
        _txb(s1, clbl,  cxj+0.12, 5.90, chip1_w-0.2, 0.28, size=9,  bold=True, color=GOLD)
        _txb(s1, cval,  cxj+0.12, 6.22, chip1_w-0.2, 0.40, size=13, bold=True, color=WHITE)

    _txb(s1, f"Prepared: {today_str}   ·   MGA Ventures — Internal Use Only",
         0.5, 6.95, 12.33, 0.38, size=10,
         color=RGBColor(0x70, 0x80, 0x98), align=PP_ALIGN.CENTER)
    # Company logo — top-right corner, white backing
    _place_logo(s1, logo_bytes, left=11.3, top=0.18, size=1.7, dark_bg=True)

    # ── Slide 2 · Company Overview ────────────────────────────────────────────
    s2 = _add_slide(prs)
    s2.background.fill.solid()
    s2.background.fill.fore_color.rgb = LTBG
    _header_bar(s2, "Company Overview", subtitle=d.get("tagline", ""))

    # Left column: extended KV table
    bm_raw = d.get("business_model", "") or ""
    bm_brief = (bm_raw[:75].rsplit(" ", 1)[0] + "…") if len(bm_raw) > 75 else bm_raw
    uof2 = d.get("use_of_funds", "Not publicly disclosed.") or "Not publicly disclosed."
    partners2 = d.get("partnerships", []) or []
    partner_str2 = "  ·  ".join(str(p) for p in partners2[:3]) if partners2 else "Not publicly disclosed."
    _kv_table(s2, [
        ("Company",        company),
        ("Sector",         d.get("sector", "")),
        ("HQ",             d.get("hq", "Not publicly disclosed.")),
        ("Founded",        d.get("founded_year", "Not publicly disclosed.")),
        ("Stage",          d.get("stage", "Not publicly disclosed.")),
        ("Website",        d.get("website", "Not publicly disclosed.")),
        ("Business Model", bm_brief or "Not publicly disclosed."),
        ("Use of Funds",   uof2),
        ("Key Partners",   partner_str2),
    ], left=0.4, top=1.15, col_widths=(2.2, 4.7))

    # Right column: highlight cards
    cards = [
        ("Stage",   d.get("stage", "Not disclosed"), NAVY),
        ("Sector",  d.get("sector", ""),             RGBColor(0x2E, 0x5F, 0xA3)),
        ("Country", "India",                          GREEN),
    ]
    cx = 8.0
    for i, (lbl, val, col) in enumerate(cards):
        card = s2.shapes.add_shape(1, Inches(cx), Inches(1.2 + i * 1.55),
                                   Inches(4.9), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = col
        card.line.fill.background()
        _txb(s2, lbl.upper(), cx + 0.18, 1.2 + i * 1.55 + 0.1, 4.5, 0.38,
             size=9, bold=True, color=RGBColor(0xB8, 0xC8, 0xE0))
        _txb(s2, val, cx + 0.18, 1.2 + i * 1.55 + 0.5, 4.5, 0.65,
             size=16, bold=True, color=WHITE)
    # Company logo — below the three cards (light background, no backing needed)
    _place_logo(s2, logo_bytes, left=9.2, top=5.72, size=2.6, dark_bg=False)

    # Recent News strip across the bottom (left side only, so logo stays clear)
    recent2 = d.get("recent_news", []) or []
    if recent2:
        rn = recent2[0]
        rn_headline = rn.get("headline", "")
        rn_date     = rn.get("date", "")
        rn_src_url  = rn.get("source_url", "")
        if rn_headline:
            rn_panel = s2.shapes.add_shape(1, Inches(0.4), Inches(6.38),
                                           Inches(8.0), Inches(0.78))
            rn_panel.fill.solid()
            rn_panel.fill.fore_color.rgb = RGBColor(0xE6, 0xEE, 0xF8)
            rn_panel.line.color.rgb = NAVY
            rn_panel.line.width = Pt(0.5)
            rn_bar = s2.shapes.add_shape(1, Inches(0.4), Inches(6.38),
                                         Inches(0.08), Inches(0.78))
            rn_bar.fill.solid()
            rn_bar.fill.fore_color.rgb = NAVY
            rn_bar.line.fill.background()
            _txb(s2, "LATEST NEWS", 0.55, 6.40, 1.6, 0.28, size=8, bold=True, color=NAVY)
            _txb(s2, rn_headline, 2.1, 6.40, 5.7, 0.38, size=12, bold=True, color=DARK)
            if rn_src_url:
                rn_outlet = re.sub(r"https?://(?:www\.)?([^/]+).*", r"\1", rn_src_url)
                _txb(s2, f"{rn_outlet}  {rn_date}".strip(), 2.1, 6.78, 6.1, 0.3,
                     size=10, italic=True, color=GRAY)

    # ── Slide 3 · What It Does ────────────────────────────────────────────────
    s3 = _add_slide(prs)
    s3.background.fill.solid()
    s3.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    _header_bar(s3, "What It Does")

    # Left column: Problem / Solution / Business Model panels (7.5" wide, 0.45" gap before right column)
    LEFT_W3 = 7.5
    sections_3 = [
        ("Problem",        "problem",        RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xBB, 0x22, 0x22)),
        ("Solution",       "solution",       RGBColor(0xF0, 0xFF, 0xF4), RGBColor(0x17, 0x6A, 0x27)),
        ("Business Model", "business_model", RGBColor(0xF0, 0xF5, 0xFF), NAVY),
    ]
    y3 = 1.05
    for lbl, key, bg, accent in sections_3:
        content = d.get(key, "Not publicly disclosed.")
        panel = s3.shapes.add_shape(1, Inches(0.33), Inches(y3),
                                    Inches(LEFT_W3), Inches(1.55))
        panel.fill.solid()
        panel.fill.fore_color.rgb = bg
        panel.line.color.rgb = accent
        panel.line.width = Pt(1)
        bar3 = s3.shapes.add_shape(1, Inches(0.33), Inches(y3),
                                   Inches(0.12), Inches(1.55))
        bar3.fill.solid()
        bar3.fill.fore_color.rgb = accent
        bar3.line.fill.background()
        _txb(s3, lbl.upper(), 0.58, y3 + 0.1, 3.0, 0.38,
             size=14, bold=True, color=accent)
        content_fit = (content[:240] + "…") if len(content) > 240 else content
        _txb(s3, content_fit, 0.58, y3 + 0.48, LEFT_W3 - 0.4, 1.05, size=13, color=DARK)
        y3 += 1.7

    # Key Differentiators — gold panel below P/S/BM (left column)
    key_diff = d.get("key_differentiators") or []
    if key_diff:
        diff_y = y3 + 0.06
        diff_h = max(7.3 - diff_y, 0.65)
        diff_panel = s3.shapes.add_shape(1, Inches(0.33), Inches(diff_y),
                                         Inches(LEFT_W3), Inches(diff_h))
        diff_panel.fill.solid()
        diff_panel.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE5)
        diff_panel.line.color.rgb = GOLD
        diff_panel.line.width = Pt(1.2)
        bar_diff = s3.shapes.add_shape(1, Inches(0.33), Inches(diff_y),
                                       Inches(0.12), Inches(diff_h))
        bar_diff.fill.solid()
        bar_diff.fill.fore_color.rgb = GOLD
        bar_diff.line.fill.background()
        _txb(s3, "KEY DIFFERENTIATORS", 0.58, diff_y + 0.09, 4.5, 0.38,
             size=12, bold=True, color=GOLD)
        diff_items = "   ·   ".join(key_diff[:4])
        _txb(s3, diff_items, 0.58, diff_y + 0.46, LEFT_W3 - 0.4, diff_h - 0.55,
             size=12, color=DARK)

    # Right column: Partnerships & Awards (4.9" wide; starts at 8.3 for a clear gap)
    RX3, RW3 = 8.3, 4.7
    rp3 = s3.shapes.add_shape(1, Inches(RX3), Inches(1.05),
                               Inches(RW3), Inches(6.3))
    rp3.fill.solid()
    rp3.fill.fore_color.rgb = RGBColor(0xEE, 0xF4, 0xFF)
    rp3.line.color.rgb = NAVY
    rp3.line.width = Pt(0.75)
    # Navy top bar for right column
    rp3_top = s3.shapes.add_shape(1, Inches(RX3), Inches(1.05), Inches(RW3), Inches(0.08))
    rp3_top.fill.solid()
    rp3_top.fill.fore_color.rgb = NAVY
    rp3_top.line.fill.background()

    _txb(s3, "PARTNERSHIPS & CUSTOMERS", RX3+0.15, 1.1, RW3-0.25, 0.38,
         size=13, bold=True, color=NAVY)
    partners3 = d.get("partnerships", []) or []
    if partners3:
        p3y = 1.55
        for pt in partners3[:5]:
            p3chip = s3.shapes.add_shape(1, Inches(RX3+0.15), Inches(p3y),
                                         Inches(RW3-0.3), Inches(0.5))
            p3chip.fill.solid()
            p3chip.fill.fore_color.rgb = WHITE
            p3chip.line.color.rgb = RGBColor(0xB0, 0xC4, 0xDC)
            p3chip.line.width = Pt(0.5)
            _txb(s3, f"  {pt}", RX3+0.25, p3y+0.08, RW3-0.5, 0.36, size=12, color=DARK)
            p3y += 0.62
    else:
        _txb(s3, "Not publicly disclosed in sources.", RX3+0.15, 1.55, RW3-0.3, 0.5,
             size=11, italic=True, color=GRAY)
        p3y = 2.2

    # Awards & Recognition below partnerships
    awards3 = d.get("awards_recognition", []) or []
    aw3_y = max(p3y + 0.25, 4.4) if partners3 else 2.3
    _txb(s3, "AWARDS & RECOGNITION", RX3+0.15, aw3_y, RW3-0.25, 0.38,
         size=13, bold=True, color=GOLD)
    if awards3:
        aw3_y += 0.42
        for aw in awards3[:4]:
            _txb(s3, f"  ★  {aw}", RX3+0.15, aw3_y, RW3-0.3, 0.42, size=11, color=DARK)
            aw3_y += 0.48
    else:
        _txb(s3, "Not publicly disclosed.", RX3+0.15, aw3_y+0.42, RW3-0.3, 0.4,
             size=11, italic=True, color=GRAY)

    # ── Slide 4 · Products ────────────────────────────────────────────────────
    s4 = _add_slide(prs)
    s4.background.fill.solid()
    s4.background.fill.fore_color.rgb = WHITE
    _header_bar(s4, "Products & Services")
    raw_products = d.get("products") or ["Not publicly disclosed."]
    # Normalise to list of {name, desc} regardless of whether Groq returned strings or dicts
    prod_items = []
    for p in raw_products:
        if isinstance(p, dict):
            prod_items.append({"name": p.get("name", ""), "desc": p.get("description", "")})
        else:
            prod_items.append({"name": str(p), "desc": ""})

    # Lay out products in 2-column cards
    if len(prod_items) >= 2:
        col_w, col_gap = 5.9, 0.5
        for i, prod in enumerate(prod_items[:6]):
            col = i % 2
            row = i // 2
            px = 0.4 + col * (col_w + col_gap)
            py = 1.2 + row * 2.0
            card4 = s4.shapes.add_shape(1, Inches(px), Inches(py),
                                        Inches(col_w), Inches(1.85))
            card4.fill.solid()
            card4.fill.fore_color.rgb = LTBG
            card4.line.color.rgb = RGBColor(0xC0, 0xCC, 0xDD)
            card4.line.width = Pt(0.75)
            bullet = s4.shapes.add_shape(1, Inches(px), Inches(py),
                                         Inches(0.1), Inches(1.85))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = GOLD
            bullet.line.fill.background()
            # Product name — bold title
            _txb(s4, prod["name"], px + 0.2, py + 0.15, col_w - 0.35, 0.45,
                 size=16, bold=True, color=NAVY)
            # Description — body text below
            if prod["desc"]:
                _txb(s4, prod["desc"], px + 0.2, py + 0.62, col_w - 0.35, 1.12,
                     size=13, color=DARK)
            else:
                _txb(s4, prod["name"], px + 0.2, py + 0.3, col_w - 0.35, 1.35,
                     size=15, color=DARK)
    else:
        names = [p["name"] if not p["desc"] else f'{p["name"]}: {p["desc"]}' for p in prod_items]
        _bullet_list(s4, names, left=0.5, top=1.2, width=12.33, height=5.8, size=18)

    # ── Slide 5 · Management Team ─────────────────────────────────────────────
    s5 = _add_slide(prs)
    s5.background.fill.solid()
    s5.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    _header_bar(s5, "Management Team")
    founders = d.get("founders") or []
    if not founders:
        _txb(s5, "Not publicly disclosed.", 0.5, 1.8, 12.33, 0.6, size=18, color=GRAY)
    else:
        cols   = 2 if len(founders) > 2 else 1
        card_w = 6.1 if cols == 2 else 12.6
        card_h = 2.6
        gap    = 0.3
        for i, f in enumerate(founders[:4]):
            col = i % cols
            row = i // cols
            fx  = 0.35 + col * (card_w + gap)
            fy  = 1.1  + row * (card_h + gap)
            fc  = s5.shapes.add_shape(1, Inches(fx), Inches(fy), Inches(card_w), Inches(card_h))
            fc.fill.solid()
            fc.fill.fore_color.rgb = WHITE
            fc.line.color.rgb = RGBColor(0xC0, 0xCC, 0xDD)
            fc.line.width = Pt(0.75)
            top_bar = s5.shapes.add_shape(1, Inches(fx), Inches(fy), Inches(card_w), Inches(0.1))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = GOLD
            top_bar.line.fill.background()
            _txb(s5, f.get("name",       ""), fx+0.2, fy+0.2,  card_w-0.35, 0.5,  size=20, bold=True,   color=NAVY)
            _txb(s5, f.get("title",      ""), fx+0.2, fy+0.72, card_w-0.35, 0.38, size=14, italic=True,  color=GRAY)
            _txb(s5, f.get("background", ""), fx+0.2, fy+1.12, card_w-0.35, 1.1,  size=14, color=DARK)
            li_url = f.get("linkedin_url", "")
            if li_url:
                _txb(s5, f"in  {li_url}", fx+0.2, fy+2.2, card_w-0.35, 0.32,
                     size=10, italic=True, color=RGBColor(0x00, 0x66, 0xCC))

        # Team Assessment strip below founder cards
        rows5  = (min(len(founders), 4) + cols - 1) // cols
        cards5_bottom = 1.1 + rows5 * (card_h + gap)
        team_strength = d.get("team_strength", "") or ""
        awards5 = d.get("awards_recognition", []) or []
        if cards5_bottom + 1.1 < 7.35 and (team_strength or awards5):
            ts_y = cards5_bottom + 0.15
            ts_h = min(7.3 - ts_y, 1.6)
            ts_panel = s5.shapes.add_shape(1, Inches(0.35), Inches(ts_y),
                                           Inches(12.63), Inches(ts_h))
            ts_panel.fill.solid()
            ts_panel.fill.fore_color.rgb = LTBG
            ts_panel.line.color.rgb = RGBColor(0xC0, 0xCC, 0xDD)
            ts_panel.line.width = Pt(0.5)
            ts_bar = s5.shapes.add_shape(1, Inches(0.35), Inches(ts_y),
                                         Inches(0.1), Inches(ts_h))
            ts_bar.fill.solid()
            ts_bar.fill.fore_color.rgb = NAVY
            ts_bar.line.fill.background()
            if team_strength and "not publicly" not in team_strength.lower():
                _txb(s5, "TEAM ASSESSMENT", 0.6, ts_y+0.1, 3.2, 0.35,
                     size=12, bold=True, color=NAVY)
                _txb(s5, team_strength, 0.6, ts_y+0.46, 7.5, ts_h-0.55,
                     size=13, color=DARK)
            if awards5:
                aw5_x = 8.2
                _txb(s5, "AWARDS & RECOGNITION", aw5_x, ts_y+0.1, 4.5, 0.35,
                     size=12, bold=True, color=GOLD)
                aw5_txt = "\n".join(f"  ★  {a}" for a in awards5[:3])
                _txb(s5, aw5_txt, aw5_x, ts_y+0.46, 4.5, ts_h-0.55, size=12, color=DARK)

    # ── Slide 6 · Market Opportunity ─────────────────────────────────────────
    s6 = _add_slide(prs)
    s6.background.fill.solid()
    s6.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    _header_bar(s6, "Market Opportunity")

    mkt = d.get("market_size") or {}
    tam = mkt.get("tam", {}) if isinstance(mkt.get("tam"), dict) else {"value": str(mkt.get("tam","Not publicly disclosed.")), "source_url": ""}
    sam = mkt.get("sam", {}) if isinstance(mkt.get("sam"), dict) else {"value": str(mkt.get("sam","Not publicly disclosed.")), "source_url": ""}
    som = mkt.get("som", {}) if isinstance(mkt.get("som"), dict) else {"value": str(mkt.get("som","Not publicly disclosed.")), "source_url": ""}
    mkt_context = mkt.get("context", d.get("market_opportunity", "Not publicly disclosed."))

    # Left: TAM / SAM / SOM attribution table
    _txb(s6, "Market Size  —  Sourced Figures", 0.35, 1.1, 7.8, 0.45,
         size=16, bold=True, color=NAVY)
    tam_sam_rows = [
        ("TAM", tam.get("value","Not publicly disclosed."), tam.get("source_url","")),
        ("SAM", sam.get("value","Not publicly disclosed."), sam.get("source_url","")),
        ("SOM", som.get("value","Not publicly disclosed."), som.get("source_url","")),
    ]
    for j, (lbl, val, src_url) in enumerate(tam_sam_rows):
        ry = 1.6 + j * 1.5
        # Label pill
        pill = s6.shapes.add_shape(1, Inches(0.35), Inches(ry), Inches(1.1), Inches(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = NAVY
        pill.line.fill.background()
        _txb(s6, lbl, 0.37, ry+0.1, 1.05, 0.38, size=16, bold=True, color=GOLD)
        # Value
        _txb(s6, val, 1.55, ry, 6.4, 0.55, size=16, bold=True, color=DARK)
        # Source attribution (smaller, italic)
        if src_url and "not publicly" not in src_url.lower():
            outlet = re.sub(r"https?://(?:www\.)?([^/]+).*", r"\1", src_url)
            _txb(s6, f"Source: {outlet}", 1.55, ry+0.58, 6.4, 0.35,
                 size=12, italic=True, color=GRAY)

    # Context paragraph below the table
    mkt_context_fit = (mkt_context[:220] + "…") if len(mkt_context) > 220 else mkt_context
    _txb(s6, mkt_context_fit, 0.35, 6.3, 7.8, 1.0, size=13, color=DARK)

    # Right: donut chart (when numeric data available) + competitive landscape below
    mkt_png = _market_donut(mkt)
    comp6_top = 1.1
    if mkt_png:
        s6.shapes.add_picture(io.BytesIO(mkt_png), Inches(8.1), Inches(1.1), width=Inches(4.9))
        comp6_top = 4.5   # competitive landscape starts below the donut

    # Competitive Landscape panel (right column, below donut or full-height if no donut)
    comp6 = d.get("competitive_landscape", []) or []
    comp6_h = max(7.25 - comp6_top, 1.0)
    cl6_panel = s6.shapes.add_shape(1, Inches(8.1), Inches(comp6_top),
                                    Inches(4.9), Inches(comp6_h))
    cl6_panel.fill.solid()
    cl6_panel.fill.fore_color.rgb = RGBColor(0xEE, 0xF4, 0xFF)
    cl6_panel.line.color.rgb = NAVY
    cl6_panel.line.width = Pt(0.75)
    cl6_top_bar = s6.shapes.add_shape(1, Inches(8.1), Inches(comp6_top),
                                      Inches(4.9), Inches(0.08))
    cl6_top_bar.fill.solid()
    cl6_top_bar.fill.fore_color.rgb = NAVY
    cl6_top_bar.line.fill.background()
    _txb(s6, "COMPETITIVE LANDSCAPE", 8.25, comp6_top + 0.1, 4.6, 0.38,
         size=13, bold=True, color=NAVY)
    if comp6:
        cl6_y = comp6_top + 0.55
        row_h6 = min(1.6, (comp6_h - 0.6) / max(len(comp6[:4]), 1))
        for cl in comp6[:4]:
            competitor = cl.get("competitor", "")
            diff6 = cl.get("differentiator", "")
            # Cap text to fit the available box height (approx 55 chars/line at size 10, 0.18" per line)
            text_h6  = max(row_h6 - 0.45, 0.4)
            max_ch6  = int(text_h6 / 0.18) * 55
            diff6_fit = (diff6[:max_ch6] + "…") if len(diff6) > max_ch6 else diff6
            cl6_row = s6.shapes.add_shape(1, Inches(8.25), Inches(cl6_y),
                                          Inches(4.6), Inches(row_h6 - 0.1))
            cl6_row.fill.solid()
            cl6_row.fill.fore_color.rgb = WHITE
            cl6_row.line.color.rgb = RGBColor(0xB8, 0xCC, 0xE4)
            cl6_row.line.width = Pt(0.5)
            _txb(s6, f"vs {competitor}", 8.4, cl6_y + 0.06, 4.3, 0.3,
                 size=11, bold=True, color=NAVY)
            _txb(s6, diff6_fit, 8.4, cl6_y + 0.36, 4.3, text_h6,
                 size=10, color=DARK)
            cl6_y += row_h6 + 0.05
    else:
        _txb(s6, "Competitive data not publicly disclosed in sources.",
             8.25, comp6_top + 0.55, 4.6, 0.5, size=12, italic=True, color=GRAY)

    # ── Slide 7 · Funding & Traction ─────────────────────────────────────────
    s7 = _add_slide(prs)
    s7.background.fill.solid()
    s7.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    _header_bar(s7, "Funding History & Traction Metrics")

    rounds  = d.get("funding_rounds")  or []
    metrics = d.get("traction_metrics") or []

    # Left half: funding table
    if rounds:
        fund_rows = [["Round", "Amount", "Investors", "Date"]]
        for r in rounds:
            fund_rows.append([r.get("round",""), r.get("amount",""),
                               r.get("investors",""), r.get("date","")])
        n   = len(fund_rows)
        tbl = s7.shapes.add_table(
            n, 4, Inches(0.35), Inches(1.15), Inches(6.5), Inches(n * 0.62)
        ).table
        for ci, cw in enumerate([1.1, 1.5, 2.8, 1.1]):
            tbl.columns[ci].width = Inches(cw)
        for ri, row in enumerate(fund_rows):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.size  = Pt(14)
                p.font.bold  = (ri == 0)
                p.font.color.rgb = WHITE if ri == 0 else DARK
                if ri == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                elif ri % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LTBG
        # Per-round source footnotes directly below the table
        fn_y = 1.15 + n * 0.62 + 0.1
        src_lines = []
        for r in rounds:
            su = r.get("source_url", "")
            if su and "not publicly" not in su.lower():
                outlet = re.sub(r"https?://(?:www\.)?([^/]+).*", r"\1", su)
                src_lines.append(f"  • {r.get('round', '?')}: {outlet} — {su}")
        if src_lines:
            note_text = "Funding sources:\n" + "\n".join(src_lines)
            note_h = 0.28 * (len(src_lines) + 1) + 0.1
            _txb(s7, note_text, 0.35, fn_y, 6.5, note_h,
                 size=10, italic=True, color=GRAY)
            fn_y += note_h + 0.1

        # Funding chart below footnotes
        fund_png = _funding_chart(rounds)
        if fund_png:
            if fn_y + 2.0 < 7.4:
                s7.shapes.add_picture(io.BytesIO(fund_png),
                                      Inches(0.35), Inches(fn_y), width=Inches(6.5))
    else:
        _txb(s7, "No funding data found in public sources.",
             0.35, 1.5, 6.5, 0.6, size=16, color=GRAY, italic=True)

    # Right half: traction chart or metric cards
    traction_png = _traction_chart(metrics)
    if traction_png:
        s7.shapes.add_picture(io.BytesIO(traction_png),
                              Inches(7.0), Inches(1.15), width=Inches(6.0))
    elif metrics:
        _txb(s7, "Traction Metrics", 7.0, 1.15, 6.0, 0.48, size=18, bold=True, color=NAVY)
        for j, m in enumerate(metrics[:4]):   # cap at 4 to leave room for Use-of-Funds strip
            my = 1.72 + j * 1.1
            card7 = s7.shapes.add_shape(1, Inches(7.0), Inches(my), Inches(6.0), Inches(0.95))
            card7.fill.solid()
            card7.fill.fore_color.rgb = WHITE
            card7.line.color.rgb = RGBColor(0xCC, 0xD8, 0xEE)
            card7.line.width = Pt(0.75)
            _txb(s7, m.get("metric",""), 7.15, my+0.08, 3.2, 0.42, size=15, bold=True, color=NAVY)
            _txb(s7, m.get("value", ""), 10.3, my+0.08, 2.55, 0.42, size=15, bold=True, color=GREEN)
            # Source attribution for this metric
            su = m.get("source_url", "")
            if su:
                outlet = re.sub(r"https?://(?:www\.)?([^/]+).*", r"\1", su)
                _txb(s7, f"Source: {outlet}", 7.15, my+0.54, 5.7, 0.3,
                     size=9, italic=True, color=GRAY)
    else:
        _txb(s7, "No traction metrics found in public sources.",
             7.0, 1.5, 6.0, 0.6, size=16, color=GRAY, italic=True)

    # Use of Funds — full-width navy strip at the very bottom of slide 7
    uof7 = d.get("use_of_funds", "") or ""
    if uof7 and "not publicly" not in uof7.lower():
        uof7_y = 6.55
        uof7_panel = s7.shapes.add_shape(1, Inches(0.35), Inches(uof7_y),
                                         Inches(12.63), Inches(0.78))
        uof7_panel.fill.solid()
        uof7_panel.fill.fore_color.rgb = NAVY
        uof7_panel.line.fill.background()
        _txb(s7, f"USE OF FUNDS:  {uof7}", 0.52, uof7_y + 0.1, 12.2, 0.58,
             size=13, color=RGBColor(0xD8, 0xE4, 0xF8))

    # ── Slide 8 · Thesis Fit ──────────────────────────────────────────────────
    s8 = _add_slide(prs)
    s8.background.fill.solid()
    s8.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    _header_bar(
        s8, "MGA Thesis Fit Assessment",
        subtitle="Stage: Seed-Series B   |   Sectors: Consumer / Fintech / SaaS   |   Ticket: up to $2M",
    )
    criteria = [
        ("Investment Stage", d.get("thesis_fit_stage",  "Not publicly disclosed.")),
        ("Sector Fit",       d.get("thesis_fit_sector", "Not publicly disclosed.")),
        ("Ticket Size",      d.get("thesis_fit_ticket", "Not publicly disclosed.")),
    ]

    # Three large cards across the top half
    card8_w = 3.9
    for i, (lbl, txt) in enumerate(criteria):
        cx8 = 0.35 + i * (card8_w + 0.22)
        cy8 = 1.2
        bg8 = s8.shapes.add_shape(1, Inches(cx8), Inches(cy8), Inches(card8_w), Inches(2.9))
        bg8.fill.solid()
        bg8.fill.fore_color.rgb = _fit_color(txt)
        bg8.line.fill.background()
        _txb(s8, lbl.upper(), cx8+0.2, cy8+0.18, card8_w-0.3, 0.5,
             size=13, bold=True, color=RGBColor(0xE8, 0xF0, 0xFF))
        verdict = "FIT" if _fit_score_val(txt) == 3 else ("PARTIAL" if _fit_score_val(txt) == 2 else "OUT OF SCOPE")
        _txb(s8, verdict, cx8+0.2, cy8+0.72, card8_w-0.3, 0.85,
             size=28, bold=True, color=WHITE)
        _txb(s8, txt, cx8+0.15, cy8+1.58, card8_w-0.25, 1.18, size=13, color=WHITE)

    # Bottom half — left 7.5": Investment Highlights stacked above Red Flags
    #               right 5.0": Gauge scorecard chart (restored to original size)
    rf_y = 4.3

    # Investment Highlights (top-left of bottom section)
    inv_highlights = d.get("investment_highlights", []) or []
    _txb(s8, "INVESTMENT HIGHLIGHTS", 0.35, rf_y, 7.5, 0.38,
         size=13, bold=True, color=GREEN)
    if inv_highlights:
        _bullet_list(s8, inv_highlights[:3], left=0.35, top=rf_y + 0.44,
                     width=7.5, height=1.3, size=13)
    else:
        _txb(s8, "See thesis fit assessment above.",
             0.35, rf_y + 0.44, 7.5, 0.4, size=13, italic=True, color=GRAY)

    # Red Flags (bottom-left of bottom section, directly below highlights)
    flags_y = rf_y + 1.85
    red_flags = d.get("red_flags") or []
    _txb(s8, "RED FLAGS", 0.35, flags_y, 7.5, 0.38, size=13, bold=True, color=RED_C)
    if red_flags:
        _bullet_list(s8, red_flags[:3], left=0.35, top=flags_y + 0.44,
                     width=7.5, height=1.1, size=13)
    else:
        _txb(s8, "No major red flags identified from public sources.",
             0.35, flags_y + 0.44, 7.5, 0.4, size=13, italic=True, color=GRAY)

    # Gauge chart — right column, original 5" width (restored)
    gauge_png = _thesis_gauge(criteria)
    if gauge_png:
        s8.shapes.add_picture(io.BytesIO(gauge_png), Inches(8.0), Inches(rf_y - 0.1),
                              width=Inches(5.0))

    # ── Slide 9 · Sources ─────────────────────────────────────────────────────
    s9 = _add_slide(prs)
    s9.background.fill.solid()
    s9.background.fill.fore_color.rgb = WHITE
    _header_bar(s9, "Sources & Attribution")

    def _is_credible_url(u: str) -> bool:
        dom = re.sub(r"https?://(?:www\.)?", "", u).split("/")[0]
        return any(dom == c or dom.endswith("." + c) for c in _CREDIBLE_DOMAINS)

    def _outlet_from_url(u: str) -> str:
        return re.sub(r"https?://(?:www\.)?([^/]+).*", r"\1", u)

    # ── Step 1: start with Groq's sources[] array ─────────────────────────────
    raw_sources = d.get("sources") or []
    seen_urls: set = set()
    all_sources: list = []
    for s in raw_sources:
        if not isinstance(s, dict):
            continue
        u = s.get("url", "")
        if u and _is_credible_url(u) and u not in seen_urls:
            seen_urls.add(u)
            all_sources.append({
                "url": u,
                "outlet": s.get("outlet", "") or _outlet_from_url(u),
                "date": s.get("date", ""),
                "claim": s.get("claim", s.get("supports", "")),
            })

    # ── Step 2: auto-harvest source_url from every structured field ───────────
    def _add_auto(url: str, claim: str, date: str = "") -> None:
        if not url or "not publicly" in url.lower():
            return
        if url in seen_urls or not _is_credible_url(url):
            return
        seen_urls.add(url)
        all_sources.append({
            "url": url,
            "outlet": _outlet_from_url(url),
            "date": date,
            "claim": claim,
        })

    # Funding rounds
    for fr in (d.get("funding_rounds") or []):
        su = fr.get("source_url", "")
        rnd = fr.get("round", "")
        amt = fr.get("amount", "")
        inv = fr.get("investors", "")
        claim_fr = " — ".join(filter(None, [rnd, amt, inv]))
        _add_auto(su, claim_fr, fr.get("date", ""))

    # Traction metrics
    for tm in (d.get("traction_metrics") or []):
        su = tm.get("source_url", "")
        claim_tm = f"{tm.get('metric','')} = {tm.get('value','')}".strip(" =")
        _add_auto(su, claim_tm)

    # Market size tiers
    mkt9 = d.get("market_size") or {}
    for tier_key, tier_label in [("tam", "TAM"), ("sam", "SAM"), ("som", "SOM")]:
        tier9 = mkt9.get(tier_key, {})
        if isinstance(tier9, dict):
            su = tier9.get("source_url", "")
            _add_auto(su, f"{tier_label} — {tier9.get('value', '')}")

    # Recent news
    for rn in (d.get("recent_news") or []):
        su = rn.get("source_url", "")
        _add_auto(su, rn.get("headline", ""), rn.get("date", ""))

    if not all_sources:
        _txb(s9, "No verifiable sources recorded.", 0.5, 1.8, 12.33, 0.6, size=18, color=GRAY)
    else:
        # ── Step 3: compact two-column cards — 6 rows × 2 cols = 12 entries ──
        col_w  = 6.0
        card_h = 0.9    # reduced from 1.4 to fit 12 entries
        row_step = 1.0  # reduced from 1.55
        for i, src in enumerate(all_sources[:12]):
            url    = src.get("url", "")
            outlet = src.get("outlet", "")
            date   = src.get("date", "")
            claim  = src.get("claim", "")

            col = i % 2
            row = i // 2
            sx  = 0.35 + col * (col_w + 0.35)
            sy  = 1.15 + row * row_step

            if sy + row_step > 7.35:
                break

            sc = s9.shapes.add_shape(1, Inches(sx), Inches(sy),
                                     Inches(col_w), Inches(card_h))
            sc.fill.solid()
            sc.fill.fore_color.rgb = LTBG
            sc.line.color.rgb = RGBColor(0xC8, 0xD4, 0xE8)
            sc.line.width = Pt(0.5)
            ab = s9.shapes.add_shape(1, Inches(sx), Inches(sy),
                                     Inches(0.08), Inches(card_h))
            ab.fill.solid()
            ab.fill.fore_color.rgb = GOLD
            ab.line.fill.background()

            header_line = f"{outlet}  {date}".strip()
            # Number + outlet clickable — links to the source
            _txb_link(s9, f"{i+1}.  {header_line}", url,
                      sx+0.15, sy+0.05, col_w-0.2, 0.26, size=11, bold=True)
            if claim:
                claim_short = (claim[:90] + "…") if len(claim) > 90 else claim
                _txb(s9, f'"{claim_short}"',
                     sx+0.15, sy+0.32, col_w-0.2, 0.30, size=10, italic=True, color=DARK)
            # URL also clickable — same target, full URL as display text
            url_display = (url[:70] + "…") if len(url) > 70 else url
            _txb_link(s9, url_display, url,
                      sx+0.15, sy+0.63, col_w-0.2, 0.22, size=9)

    # ── Save (add _v2/_v3 if today's file is already open) ───────────────────
    safe = re.sub(r"[^\w\-]", "_", company)
    out  = BASE_DIR / f"{safe}_{today_str}.pptx"
    if out.exists():
        for ver in range(2, 30):
            candidate = BASE_DIR / f"{safe}_{today_str}_v{ver}.pptx"
            if not candidate.exists():
                out = candidate
                break
    prs.save(out)
    print(f"[Deck] Saved -> {out.name}")
    return out


# ─────────────────────────────────────────────────────────────────────────────
# 7. Gmail send
# ─────────────────────────────────────────────────────────────────────────────

def _gmail_service():
    """
    Authenticate with Gmail API using OAuth2.
    On first run: opens a browser window for one-time consent.
    On subsequent runs: uses the saved token.json silently.
    """
    creds = None
    if TOKEN_PATH.exists():
        creds = Credentials.from_authorized_user_file(str(TOKEN_PATH), GMAIL_SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(GoogleRequest())
        else:
            if not CREDS_PATH.exists():
                raise RuntimeError(
                    f"{CREDS_PATH.name} not found. "
                    "Gmail OAuth requires credentials.json from the Google Cloud Console. "
                    "On Streamlit Cloud, set the GMAIL_TOKEN_JSON secret instead."
                )
            flow  = InstalledAppFlow.from_client_secrets_file(str(CREDS_PATH), GMAIL_SCOPES)
            creds = flow.run_local_server(port=0)
        TOKEN_PATH.write_text(creds.to_json())
    return build("gmail", "v1", credentials=creds)


def email_deck(pptx_path: Path, company: str, summary: str, today_str: str):
    """Send the PowerPoint as a Gmail attachment."""
    if not RECIPIENT_EMAIL:
        print("[WARN] RECIPIENT_EMAIL not set in .env — skipping email.")
        return

    print(f"[Gmail] Sending to {RECIPIENT_EMAIL} …")
    svc = _gmail_service()

    msg             = MIMEMultipart()
    msg["to"]       = RECIPIENT_EMAIL
    msg["subject"]  = f"MGA Deal Scout — {company} — {today_str}"
    msg.attach(MIMEText(summary, "plain"))

    with open(pptx_path, "rb") as fh:
        part = MIMEBase("application", "octet-stream")
        part.set_payload(fh.read())
    encoders.encode_base64(part)
    part.add_header("Content-Disposition",
                    f'attachment; filename="{pptx_path.name}"')
    msg.attach(part)

    raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
    svc.users().messages().send(userId="me", body={"raw": raw}).execute()
    print(f"[Gmail] Sent.")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    today_str = datetime.date.today().strftime("%Y-%m-%d")
    print(f"\n{'=' * 62}")
    print(f"  MGA Deal Scout   ·   {today_str}")
    print(f"{'=' * 62}\n")

    # Check required env vars before doing anything
    missing = [k for k in ("SERPER_API_KEY", "GROQ_API_KEY")
               if not os.getenv(k)]
    if missing:
        print(
            f"[ERROR] Missing environment variables: {', '.join(missing)}\n"
            "Copy .env.example -> .env and fill in your API keys.\n"
            "See the setup walkthrough in the script comments or README."
        )
        sys.exit(1)

    # Step 0 — sector (CLI override for testing; weekly rotation in production)
    parser = argparse.ArgumentParser(description="MGA Deal Scout")
    parser.add_argument(
        "--sector",
        choices=SECTORS,
        default=None,
        help="Override weekly rotation: Consumer | Fintech | SaaS",
    )
    args = parser.parse_args()
    sector = args.sector if args.sector else sector_this_week()
    src    = f"(override)" if args.sector else f"(week rotation)"
    print(f"Sector this week: {sector}  {src}\n")

    # Step 1 — discover one candidate startup
    candidate = discover_startup(sector)
    company   = candidate["name"]
    print(f"\nSelected candidate: {company}\n")

    # Step 2 — deep research (4 CSE queries + page fetching)
    research = deep_research(company, sector, candidate["url"])

    # Step 3 — synthesise with Groq
    structured = synthesise(company, sector, research)

    # Step 3b — validate India HQ from Groq's output
    # If Groq found a non-India headquarters, abort rather than send a wrong deck.
    hq = structured.get("hq", "")
    _INDIA_CITIES = {
        "india", "bengaluru", "bangalore", "mumbai", "delhi", "hyderabad",
        "pune", "chennai", "noida", "gurugram", "gurgaon", "kolkata",
        "ahmedabad", "jaipur", "indore", "surat", "kochi", "coimbatore",
    }
    hq_lower = hq.lower()
    hq_is_india = any(city in hq_lower for city in _INDIA_CITIES)
    if hq and not hq_is_india and "not publicly" not in hq_lower:
        print(f"[WARN] Groq identified HQ as '{hq}' — not an Indian city. Aborting.")
        print("       Re-run to pick a different candidate, or use --sector to try another sector.")
        sys.exit(0)

    # Step 3c — fetch company logo (Clearbit, free, no auth)
    website = structured.get("website", "")
    logo_bytes = _fetch_logo(website)

    # Step 4 — build PowerPoint
    pptx_path = build_deck(structured, today_str, logo_bytes=logo_bytes)

    # Step 5 — email the deck
    summary = structured.get(
        "email_summary",
        f"{company} — see attached deal screening memo.",
    )
    email_deck(pptx_path, structured.get("company_name", company), summary, today_str)

    print(f"\n{'=' * 62}")
    print(f"  Done.")
    print(f"  Deck    : {pptx_path.name}")
    print(f"  Email   : {'sent to ' + RECIPIENT_EMAIL if RECIPIENT_EMAIL else 'skipped (no RECIPIENT_EMAIL)'}")
    print(f"{'=' * 62}\n")


if __name__ == "__main__":
    main()
